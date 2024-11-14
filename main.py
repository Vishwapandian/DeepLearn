import os
import re
import PyPDF2
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from pathlib import Path
from openai import OpenAI
import subprocess
from pdf2image import convert_from_path

# Import pptx library for creating PPTX files
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize OpenAI client with API key
client = OpenAI(
    api_key = ''
)

def generate_slides_and_scripts(text):
    # System prompt to set the context and expectations
    system_prompt = """You are an expert university professor and instructional designer specializing in creating engaging educational content. Your task is to create a college-level presentation that:

1. Structures complex information in a clear, logical sequence
2. Uses academic language while remaining accessible
3. Follows educational best practices for slide design:
   - Clear hierarchy of information
   - 3-5 key points per slide
   - Meaningful titles that convey main ideas
   - Bullet points that support learning objectives
4. Creates engaging narrative scripts that:
   - Expand significantly on slide content
   - Include relevant examples and case studies
   - Connect concepts to real-world applications
   - Use rhetorical questions and thought experiments
   - Maintain a conversational yet professional tone

Each slide should build upon previous concepts and contribute to a cohesive learning experience."""

    # Generate slides content and scripts using the extracted text
    prompt = f"""Create a detailed university-level presentation based on the following text. Structure the presentation to maximize student learning and engagement.

For each slide, provide:
1. A clear, conceptual title that frames the main idea
2. 3-5 carefully crafted bullet points that:
   - Present key academic concepts
   - Use precise terminology
   - Support progressive understanding
   - Include relevant data or evidence when applicable

3. A detailed script (2-3 minutes) that:
   - Provides deeper context and explanations
   - Includes specific examples and case studies
   - Connects to broader academic frameworks
   - Poses thought-provoking questions
   - Explains complex relationships between concepts
   - Uses analogies to clarify difficult ideas
   - Maintains student engagement through narrative techniques

Format the output exactly as follows:

### Slide 1:
Title: [Conceptual Title]
Bullet Points:
- [Academic point with supporting evidence]
- [Key concept with precise terminology]
- [Critical relationship or framework]
Script:
[Detailed narrative that significantly expands on the bullet points while maintaining engagement]

### Slide 2:
[Continue format...]

Source Text:
{text}"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",  # Using the latest model for best quality
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt}
        ],
        max_tokens=16384,
        temperature=0.7,  # Slightly increased for more creative and engaging scripts
    )
    
    slides_and_scripts = response.choices[0].message.content.strip()
    return slides_and_scripts

def generate_audio(script, slide_number, voice="alloy"):
    """Generate audio using OpenAI's TTS API with enhanced speaking instructions"""
    
    # Add presentation instructions to the script
    enhanced_script = f"""[Speak in an engaging, professional tone appropriate for a university lecture. Maintain a steady, clear pace with natural pauses for emphasis and comprehension. Use vocal variety to highlight key points and maintain student attention.]

{script}"""
    
    speech_file_path = Path(f"audio_{slide_number}.mp3")
    
    response = client.audio.speech.create(
        model="tts-1",  # Using HD model for better quality
        voice=voice,
        input=enhanced_script,
        speed=0.9  # Slightly slower for better comprehension
    )
    
    response.stream_to_file(str(speech_file_path))
    return str(speech_file_path)

def parse_slides_and_scripts(content):
    slides = []
    # Split content into individual slides
    slide_sections = re.split(r'\n### Slide \d+:\n', content)
    for slide_section in slide_sections[1:]:  # Skip the first empty split
        # Extract title
        title_match = re.search(r'Title:\s*(.+)', slide_section)
        title = title_match.group(1).strip() if title_match else "Untitled Slide"
        # Extract bullet points
        bullet_points = re.findall(r'- (.+)', slide_section)
        # Extract script
        script_match = re.search(r'Script:\n(.+)', slide_section, re.DOTALL)
        script = script_match.group(1).strip() if script_match else ""
        slides.append({'title': title, 'bullet_points': bullet_points, 'script': script})
    return slides

def generate_audio(script, slide_number, voice="alloy"):
    """Generate audio using OpenAI's TTS API"""
    speech_file_path = Path(f"audio_{slide_number}.mp3")
    
    response = client.audio.speech.create(
        model="tts-1",
        voice=voice,
        input=script
    )
    
    # Save the audio file
    response.stream_to_file(str(speech_file_path))
    
    return str(speech_file_path)

def export_slides_to_images(pptx_filename, output_folder='slides'):
    # Create output folder if it doesn't exist
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Define the paths
    pptx_path = os.path.abspath(pptx_filename)
    pdf_path = os.path.join(output_folder, "presentation.pdf")

    # Step 1: Convert PPTX to PDF
    command = [
        "soffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_folder,
        pptx_path
    ]
    subprocess.run(command, check=True)

    # Step 2: Convert PDF to individual PNG images
    slide_filenames = []
    images = convert_from_path(pdf_path)
    for i, image in enumerate(images):
        slide_filename = os.path.join(output_folder, f"slide_{i+1}.png")
        image.save(slide_filename, "PNG")
        slide_filenames.append(slide_filename)

    # Remove the intermediate PDF file if not needed
    os.remove(pdf_path)

    return slide_filenames

def extract_text_from_pdf(pdf_path):
    pdfReader = PyPDF2.PdfReader(open(pdf_path, 'rb'))
    text = ""
    for page in pdfReader.pages:
        text += page.extract_text()
    return text

def create_intro_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    # Set title
    title_placeholder.text = title_text
    title_tf = title_placeholder.text_frame
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Set subtitle
    subtitle_placeholder.text = subtitle_text
    subtitle_tf = subtitle_placeholder.text_frame
    subtitle_tf.paragraphs[0].font.size = Pt(28)
    subtitle_tf.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Set background color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 70, 127)  # Dark blue

def create_outro_slide(prs, thank_you_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank Slide layout
    left = Inches(0)
    top = Inches(2)
    width = prs.slide_width
    height = Inches(3)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = thank_you_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Set background color
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 70, 127)  # Dark blue

def create_presentation(slides, pptx_filename='presentation.pptx', presenter_name="DeepLearn"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    presentation_title = slides[0]['title'] if slides else "Presentation"
    create_intro_slide(prs, presentation_title, f"Presented by {presenter_name}")

    for slide_content in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using Title and Content layout
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        # Set the title with improved spacing and size
        title_placeholder.text = slide_content['title']
        title_tf = title_placeholder.text_frame
        title_tf.paragraphs[0].font.size = Pt(40)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 70, 127)
        title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT  # Align to left for consistency

        # Add bullet points with more controlled spacing and font adjustments
        tf = content_placeholder.text_frame
        tf.clear()  # Clear any existing content
        for bullet_point in slide_content['bullet_points']:
            p = tf.add_paragraph()
            p.text = bullet_point
            p.level = 0
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.alignment = PP_ALIGN.LEFT

        # Set background color to improve visual contrast
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background for better contrast

    # Create outro slide
    create_outro_slide(prs, "Thank You!")
    prs.save(pptx_filename)

def create_video(slide_filenames, audio_filenames, output_filename="presentation.mp4"):
    clips = []
    for slide_filename, audio_filename in zip(slide_filenames, audio_filenames):
        # Get audio duration
        audio_clip = AudioFileClip(audio_filename)
        duration = audio_clip.duration

        # Create ImageClip with duration equal to audio duration
        image_clip = ImageClip(slide_filename).set_duration(duration)

        # Set audio
        image_clip = image_clip.set_audio(audio_clip)

        clips.append(image_clip)

    # Concatenate clips
    final_clip = concatenate_videoclips(clips, method="compose")

    # Write the video file
    final_clip.write_videofile(output_filename, fps=24)

    # Close clips to release resources
    final_clip.close()
    for clip in clips:
        clip.close()

# Updated main function
def main(pdf_path, voice="alloy"):
    # Step 1: Extract text from PDF
    print("Extracting text from PDF...")
    text = extract_text_from_pdf(pdf_path)

    # Step 2: Generate slides content and scripts
    print("Generating slides content and scripts...")
    slides_and_scripts = generate_slides_and_scripts(text)

    # Parse the slides content and scripts
    slides = parse_slides_and_scripts(slides_and_scripts)

    # List to track temp files
    temp_files = []

    # Step 3: Create PPTX presentation
    print("Creating PPTX presentation...")
    pptx_filename = 'presentation.pptx'
    create_presentation(slides, pptx_filename, presenter_name="DeepLearn")
    temp_files.append(pptx_filename)  # Track the PPTX file

    # Step 4: Export slides to images
    print("Exporting slides to images...")
    slide_filenames = export_slides_to_images(pptx_filename)
    temp_files.extend(slide_filenames)  # Track all slide image files

    # Ensure the number of slide images matches the number of slides plus intro and outro
    expected_slides = len(slides) + 2  # Intro and Outro slides
    if len(slide_filenames) != expected_slides:
        print("Error: Number of slide images does not match expected number of slides.")
        return

    audio_filenames = []

    # Generate audio for intro slide
    intro_script = f"Welcome to this presentation on {slides[0]['title']}. Presented by DeepLearn."
    audio_filename = generate_audio(intro_script, 1, voice=voice)
    audio_filenames.append(audio_filename)
    temp_files.append(audio_filename)

    # Generate audio for main slides using the scripts from the slides
    for idx, slide in enumerate(slides):
        print(f"Processing Slide {idx+1}: {slide['title']}")

        # Use the script generated earlier
        script = slide['script']

        # Step 5: Generate audio
        audio_filename = generate_audio(script, idx+2, voice=voice)  # idx+2 because intro slide is first
        audio_filenames.append(audio_filename)
        temp_files.append(audio_filename)

    # Generate audio for outro slide
    outro_script = "Thank you for watching this presentation."
    audio_filename = generate_audio(outro_script, len(slides)+2, voice=voice)
    audio_filenames.append(audio_filename)
    temp_files.append(audio_filename)

    # Step 6: Create video
    print("Creating video presentation...")
    output_filename = "presentation.mp4"
    create_video(slide_filenames, audio_filenames, output_filename)
    print("Video presentation created successfully!")

    # Clean up temporary files, keeping only the final MP4
    print("Cleaning up temporary files...")
    for temp_file in temp_files:
        try:
            os.remove(temp_file)
        except OSError as e:
            print(f"Error deleting {temp_file}: {e}")

    print("Temporary files deleted. Only the final MP4 file remains.")

if __name__ == '__main__':
    pdf_path = 'input.pdf'  # Replace with your PDF file path
    voice = 'alloy'  # Can be: alloy, echo, fable, onyx, nova, or shimmer
    main(pdf_path, voice)
